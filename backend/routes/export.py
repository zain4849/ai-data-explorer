"""Export endpoints: CSV, Excel, PDF, PNG, PPTX."""

import io
import json

import pandas as pd
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from ..auth.dependencies import get_current_user
from ..connectors import get_connector
from ..connectors.types import ConnectionConfig
from ..crypto import decrypt
from ..db import db_manager
from ..models.base import get_session
from ..models.connection import DataConnection
from ..models.user import User
from ..sql_validator import ensure_limit, validate_sql
from sqlalchemy.orm import Session

router = APIRouter(prefix="/export", tags=["export"])


class ExportRequest(BaseModel):
    sql: str | None = None
    connection_id: str | None = None
    chart_html: str | None = None
    data: list[dict] | None = None


def _get_dataframe(body: ExportRequest, user: User, session: Session) -> pd.DataFrame:
    """Execute the SQL to produce a DataFrame for export."""
    if body.data:
        return pd.DataFrame(body.data)

    if not body.sql:
        raise HTTPException(status_code=400, detail="sql or data is required")

    validate_sql(body.sql)
    sql = ensure_limit(body.sql, max_limit=10000)

    if body.connection_id:
        conn = session.get(DataConnection, body.connection_id)
        if conn is None or conn.user_id != user.id:
            raise HTTPException(status_code=404, detail="Connection not found")
        config = ConnectionConfig(**json.loads(decrypt(conn.config_encrypted)))
        connector = get_connector(conn.db_type, config)
        try:
            return connector.execute_query(sql)
        finally:
            connector.close()
    else:
        return db_manager.query(user.id, sql)


@router.post("/csv")
def export_csv(body: ExportRequest, user: User = Depends(get_current_user), session: Session = Depends(get_session)):
    df = _get_dataframe(body, user, session)
    buf = io.StringIO()
    df.to_csv(buf, index=False)
    buf.seek(0)
    return StreamingResponse(
        iter([buf.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=export.csv"},
    )


@router.post("/xlsx")
def export_xlsx(body: ExportRequest, user: User = Depends(get_current_user), session: Session = Depends(get_session)):
    df = _get_dataframe(body, user, session)
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Data")
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=export.xlsx"},
    )


@router.post("/png")
def export_png(body: ExportRequest, user: User = Depends(get_current_user), session: Session = Depends(get_session)):
    if not body.chart_html:
        raise HTTPException(status_code=400, detail="chart_html is required for PNG export")
    try:
        import plotly.io as pio
        fig = pio.from_json(body.chart_html) if body.chart_html.startswith("{") else None
        if fig is None:
            raise HTTPException(status_code=400, detail="Invalid chart data for image export")
        img_bytes = fig.to_image(format="png", engine="kaleido")
        return StreamingResponse(
            io.BytesIO(img_bytes),
            media_type="image/png",
            headers={"Content-Disposition": "attachment; filename=chart.png"},
        )
    except ImportError:
        raise HTTPException(status_code=501, detail="kaleido is required for PNG export")


@router.post("/pdf")
def export_pdf(body: ExportRequest, user: User = Depends(get_current_user), session: Session = Depends(get_session)):
    raise HTTPException(
        status_code=501,
        detail="PDF export not yet implemented. Use CSV, XLSX, or PPTX for data export.",
    )


@router.post("/pptx")
def export_pptx(body: ExportRequest, user: User = Depends(get_current_user), session: Session = Depends(get_session)):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise HTTPException(status_code=501, detail="python-pptx is required for PPTX export")

    df = _get_dataframe(body, user, session)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    title = slide.shapes.title
    if title:
        title.text = "Data Export"

    # Add table
    rows, cols = min(len(df), 20), len(df.columns)
    table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    table = table_shape.table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i in range(rows):
        for j in range(cols):
            val = df.iloc[i, j]
            table.cell(i + 1, j).text = "" if pd.isna(val) else str(val)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=export.pptx"},
    )
